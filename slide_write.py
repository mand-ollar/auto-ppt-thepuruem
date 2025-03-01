from pptx import Presentation
from get_phrases import getPhrases
from pathlib import Path
from lyrics import get_lyrics

import json
import re
import os
import keyboard


class slideWrite:
    def __init__(self,
                 ppt_file: Presentation = None,
                 layout_dict: dict = None,
                 variable_list: list = None,
                 ):
        '''
        \uAC00-\uD7A30 : 가-힣
        0-9 : 0-9
        a-zA-Z : a-zA-Z
        \s : whitespace
        '''

        self.get_lyrics = get_lyrics.getLyrics()

        self.placeholer_dict = {
            "preview": ["pdate", "title", "bible"],
            "worshiptitle": ["title_kr", "title_en"],
            "worshiplyrics": ["lyrics"],
            "prayer": ["name"],
            "sermontitle": ["preacher", "title", "bible"],
            "biblephrase": ["phrase", "loc"],
            "intro": ["adate"],
        }

        self.variable_dict = {
            "pdate": f"{variable_list[0]} 주일 예배",
            "title": f'"{variable_list[1]}"',
            "bible": variable_list[2],
            "insermon": variable_list[3],   # list
            "title_kr": variable_list[4],   # list
            "title_en": variable_list[5],   # list
            "lyrics": variable_list[6],     # list
            "sequence": variable_list[7],   # list
            "name": f"기도 | {variable_list[8]}",
            "preacher": f"말 씀 | {variable_list[9]}",
            "adate": variable_list[0],
            "version": variable_list[10],
            "loc": None,
            "phrase": None,
        } if variable_list is not None else None

        self.ppt_file = ppt_file if ppt_file is not None else None
        self.layout_dict = layout_dict if layout_dict is not None else None

        self.gp = getPhrases()

    def _bible_to_book_chap_sec(self, bibles):
        
        books = []; chaps = []; secs = []
        for bible in bibles:

            # Extract book
            current_file_dir = Path(__file__).parent
            json_dict = current_file_dir/"dict/book-dict.json"
            with open(json_dict, "r") as f:
                book_dict = json.load(f)
                f.close()

            book_list = list(book_dict.keys())
            find_status = False
            for book in book_list:
                if book in bible:
                    find_status = True
                    break

            if not find_status:
                print("Error: Cannot find the book in the Bible.")
                return None

            if ":" in bible:
                chap_sec = bible.replace(book, "").strip()
                idx = chap_sec.find(":")
                chap = chap_sec[:idx].strip()
                sec = chap_sec[idx+1:].strip()
                
            elif "장" in bible and "절" in bible:
                chap_sec = bible.replace(book, "").strip()
                chap = bible[bible.find("장") - 1]
                sec = bible[bible.find("절") - 1]

            if "-" in sec:
                idx = sec.find("-")
                
                init = int(sec[:idx])
                end = int(sec[idx+1:])
                
                sec = list(range(init, end+1))

            elif "~" in sec:
                idx = sec.find("~")
                
                init = int(sec[:idx])
                end = int(sec[idx+1:])
                
                sec = list(range(init, end+1))
                
            else:
                sec = [int(sec)]

            books.append(book)
            chaps.append(chap)
            secs.append(sec)
            
        return books, chaps, secs

    def _find_placeholder_by_name(self, placeholder_name, placeholders):
        for i in placeholders:
            if i.name == placeholder_name:
                return i

    def write_contents(self, layout_name):

        if layout_name == "allofsermon":
            
            # set layout_name
            layout_name = "sermontitle"
                
            # add slide
            slide = self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
            
            # find placeholders
            phs = slide.placeholders
            
            # get number of placeholders
            ph_cnt = 0
            for ph in phs:
                ph_cnt += 1
                
            # write contents
            for i in range(ph_cnt):
                ph = self._find_placeholder_by_name(f"Text Placeholder {i+1}", phs)
                ph.text = self.variable_dict[self.placeholer_dict[layout_name][i]]
            
            # bible -> book, chap, sec
            bibles = [self.variable_dict["bible"]] + self.variable_dict["insermon"]
            books, chaps, secs = self._bible_to_book_chap_sec(bibles)
            
            for k in range(len(bibles)):
                for i in range(len(secs[k])):
                    
                    # set layout_name
                    layout_name = "biblephrase"
                    
                    # add slide
                    slide = self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
                    
                    # find placeholders
                    phs = slide.placeholders
                    
                    secnum = secs[k][i]
                    phrase = self.gp.get_phrase(version=self.variable_dict["version"],
                                                book=books[k],
                                                chap=chaps[k],
                                                sec=secnum,
                                                )
                    
                    self.variable_dict["phrase"] = phrase
                    self.variable_dict["loc"] = f"{books[k]} {chaps[k]}:{secnum}"
                    
                    # get number of placeholders
                    ph_cnt = 0
                    for ph in phs:
                        ph_cnt += 1
                    
                    # write contents
                    for j in range(ph_cnt):
                        ph = self._find_placeholder_by_name(f"Text Placeholder {j+1}", phs)
                        ph.text = self.variable_dict[self.placeholer_dict[layout_name][j]]
                
                # set layout_name
                layout_name = "sermontitle"
                    
                # add slide
                slide = self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
                
                # find placeholders
                phs = slide.placeholders
                
                # get number of placeholders
                ph_cnt = 0
                for ph in phs:
                    ph_cnt += 1
                    
                # write contents
                for i in range(ph_cnt):
                    ph = self._find_placeholder_by_name(f"Text Placeholder {i+1}", phs)
                    ph.text = self.variable_dict[self.placeholer_dict[layout_name][i]]

        elif layout_name not in self.placeholer_dict:
            self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
        
        elif layout_name not in ["allofsermon"]:
            # add slide
            slide = self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
            
            # find placeholders
            phs = slide.placeholders
            
            # get number of placeholders
            ph_cnt = 0
            for ph in phs:
                ph_cnt += 1
            
            # write contents
            for i in range(ph_cnt):
                ph = self._find_placeholder_by_name(f"Text Placeholder {i+1}", phs)
                ph.text = self.variable_dict[self.placeholer_dict[layout_name][i]]

    def _prepare_lyrics_(self,
                        title,
                        ):
        loop = True
        while loop == True:
            title_for_bash = title.replace(" ", "\ ")
            self.get_lyrics.get_lyrics(song_title=title, artist_name="thisisnotanartistname")
            current_dir = Path(__file__).parent
            os.system(f"nano {current_dir}/lyrics/txt/{title_for_bash}.txt")
            loop = self.get_lyrics.get_lyrics(song_title=title, artist_name="thisisnotanartistname", print_it=False)

            if loop:
                print()
                print("Press <enter> to try again with text edit.")
                print("Or press <again()> if you want to choose artist again.")
                print("Or press <exit()> if you want to exit.")

                real_break = False
                key_input = "thisisnotaninput"
                while key_input == "thisisnotaninput":
                    key_input = input(": ")

                    if key_input == "":
                        continue

                    elif key_input == "again()":
                        Path(f"{current_dir}/lyrics/txt/{title}.txt").unlink()
                        os.system("clear")

                    elif key_input == "exit()":
                        Path(f"{current_dir}/lyrics/txt/{title}.txt").unlink()
                        real_break = True
                        break

                    else:
                        key_input = "thisisnotaninput"
                        print("Invalid input. Try again.")

                print()
                if real_break: break

    def write_lyrics(self,
                     lyrics,    # dict
                     sequence,  # str ex> "A B A B C"  whitespace, lowercase available
                     ):

        layout_name = "worshiplyrics"

        sequence = sequence.upper().replace(" ", "")

        # find alphabets idx
        alph_idx = []
        for i in range(len(sequence)):
            if sequence[i].isalpha():
                alph_idx.append(i)
        alph_idx.append(len(sequence))

        # segment list
        temp = []
        for i in range(len(alph_idx)-1):
            temp.append(sequence[alph_idx[i]:alph_idx[i+1]][:])

        sequence = temp

        # make lyrics dict by parts {A: list, B: list, ...}
        parts_for_ppt = {}
        for part in lyrics:
            lyric = lyrics[part].split("\n")

            lyric_segments = []

            n = len(lyric)
            for i in range(round(n/2 + 0.1)):
                if n != 2*i + 1:
                    lyric_segments.append(lyric[2 * i] + '\n' + lyric[2 * i + 1])
                else:
                    lyric_segments.append(lyric[2 * i])

            parts_for_ppt[part] = lyric_segments
        # part by part
        for seq in sequence:

            if seq == "Z":
                self.ppt_file.slides.add_slide(self.layout_dict["worshipbg"])
                continue

            # get lyrics for a slide
            lys = parts_for_ppt[seq]

            for seg in lys:

                #add slide
                slide = self.ppt_file.slides.add_slide(self.layout_dict[layout_name])
                
                # find placeholders
                phs = slide.placeholders
                
                # get number of placeholders
                ph_cnt = 0
                for ph in phs:
                    ph_cnt += 1

                # write lyrics
                for i in range(ph_cnt):
                    ph = self._find_placeholder_by_name(f"Text Placeholder {i+1}", phs)
                    ph.text = seg

    def write_worship(self,
                      type='remain',
                      title=None,
                      ):
        if type == "remain":
            for i in range(len(self.variable_dict["title_kr"])):
                title_kr = self.variable_dict["title_kr"][i]
                title_en = self.variable_dict["title_en"][i]
                lyrics_rm = self.variable_dict["lyrics"][i]
                sequence = self.variable_dict["sequence"][i]

                # add title slide
                slide = self.ppt_file.slides.add_slide(self.layout_dict["worshiptitle"])

                # find placeholders
                phs = slide.placeholders

                # write titles
                ph = self._find_placeholder_by_name(f"Text Placeholder 1", phs)
                ph.text = title_kr

                ph = self._find_placeholder_by_name(f"Text Placeholder 2", phs)
                ph.text = title_en

                # add lyrics slide
                self.write_lyrics(lyrics_rm, sequence)
        
        elif type == "normal":
            
            # add title slide
            slide = self.ppt_file.slides.add_slide(self.layout_dict["worshiptitle"])
            
            # find placeholders
            phs = slide.placeholders
            
            # write titles
            ph = self._find_placeholder_by_name(f"Text Placeholder 1", phs)
            ph.text = title
            
            ph = self._find_placeholder_by_name(f"Text Placeholder 2", phs)
            ph.text = ""
            
            # if lyrics json exists
            if Path(f"lyrics/json/{title}.json").exists():
                with open(f"lyrics/json/{title}.json", "r") as lyrics_json:
                    lyrics = json.load(lyrics_json)
            # if lyrics json doesn't exist
            else:
                print(f"Lyrics file for [{title}] doesn't exist...\n")
                print("Terminating...")
                self.get_lyrics.get_lyrics(title)
                return False
            
            # sequence setting
            sequence_list = list(lyrics.keys())
            sequence = "".join(sequence_list)
            sequence = sequence + "z"
            
            # add lyrics slide
            self.write_lyrics(lyrics, sequence)


if __name__ == "__main__":
    sW = slideWrite(ppt_file = None,
                    layout_dict = None,
                    variable_list = None,
                    )
    sW._prepare_lyrics_("그가 내 안에")
